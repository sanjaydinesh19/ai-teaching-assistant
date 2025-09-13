import os, time, random, json
from typing import List, Dict, Any, Tuple
from PyPDF2 import PdfReader
from pptx import Presentation
from openai import OpenAI

from .pdf import render_pdf
from .schemas import WorksheetRequest, WorksheetResponse, WorksheetSetResponse, WorksheetItem

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4

FILE_ROOT = os.environ.get("FILE_STORE", "/data")
OPENAI_MODEL = os.environ.get("WORKSHEET_MODEL", "gpt-4o-mini")

client = OpenAI()


def _id(prefix="ws"):
    t = str(int(time.time()))
    rnd = "".join(random.choice("abcdefghijklmnopqrstuvwxyz0123456789") for _ in range(6))
    return f"{prefix}_{t}{rnd}"


# ---- Source reading ----
def _find_path(file_id: str) -> Tuple[str, str]:
    for ext in [".pdf", ".pptx", ".ppt", ".png", ".jpg", ".jpeg", ".webp"]:
        p = os.path.join(FILE_ROOT, f"{file_id}{ext}")
        if os.path.exists(p):
            print("[DEBUG] Listing /data:", os.listdir(FILE_ROOT))
            return p, ext.lower()
    return os.path.join(FILE_ROOT, file_id), ""


def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    chunks = []
    for page in reader.pages:
        try:
            chunks.append(page.extract_text() or "")
        except Exception:
            pass
    return "\n".join(chunks).strip()


def extract_text_from_ppt(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts).strip()


def aggregate_source_text(file_ids: List[str]) -> str:
    blobs = []
    for fid in file_ids:
        path, ext = _find_path(fid)
        if not os.path.exists(path):
            blobs.append(f"[WARN] Missing file: {fid}")
            continue
        if ext == ".pdf":
            blobs.append(f"[PDF:{os.path.basename(path)}]\n" + extract_text_from_pdf(path))
        elif ext in (".pptx", ".ppt"):
            blobs.append(f"[PPT:{os.path.basename(path)}]\n" + extract_text_from_ppt(path))
        elif ext in (".png", ".jpg", ".jpeg", ".webp"):
            blobs.append(f"[IMAGE:{os.path.basename(path)}] (content requires OCR/vision; generate items aligned to grade + topic context)")
        else:
            blobs.append(f"[UNKNOWN:{os.path.basename(path)}]")
    text = "\n\n".join(blobs)
    return text[:15000]


# ---- LLM generation ----
SYS = """You are a worksheet designer for multi-grade classrooms in India.
Create classroom-ready items suitable for the requested difficulty and grades.
Prefer low-resource, practical prompts. Avoid copyrighted passages.
Return STRICT JSON only with key "items" as an array of objects:
{ "type": "mcq|short|diagram|long", "q": "question text", "options": ["A","B","C","D"]?, "answer": "correct or sample", "rubric": "optional" }
Do not include any other keys.
"""


def generate_items(context: str, req: WorksheetRequest, difficulty: str, total_q: int) -> List[Dict[str, Any]]:
    mix_hint = ""
    if req.question_mix:
        parts = [f'{k}:{v}' for k, v in req.question_mix.items()]
        mix_hint = " Question mix: " + ", ".join(parts) + "."

    grades_hint = ""
    if req.grade_bands:
        grades_hint = f" Grade bands: {', '.join(req.grade_bands)}."

    user = f"""
TARGET LANGUAGE: {req.target_language}
GRADE LEVELS: {', '.join(req.grade_bands or [])}

You are creating a classroom worksheet with DIFFICULTY = {difficulty.upper()}.
Each worksheet must be written in {req.target_language.upper()}.
Generate exactly {total_q} questions, using the SOURCE CONTEXT when possible.

Difficulty must be defined as:
- EASY → recall or recognition
- MEDIUM → application or explanation
- HARD → higher-order analysis

Rules:
- Do not simply reword the same question across difficulties
- Always prefer context-relevant content. If context is thin or missing, fallback to typical grade topics.

SOURCE CONTEXT (may be partial):
\"\"\"{context}\"\"\"

Return STRICT JSON with key "items".
"""

    resp = client.chat.completions.create(
        model=OPENAI_MODEL,
        temperature=0.7,
        messages=[
            {"role": "system", "content": SYS},
            {"role": "user", "content": user}
        ],
        response_format={"type": "json_object"}
    )
    data = json.loads(resp.choices[0].message.content)
    return data.get("items", [])


# ---- Public entry ----
def build_worksheets(req: WorksheetRequest) -> WorksheetResponse:
    context = aggregate_source_text(req.file_ids)

    if len(req.difficulty_levels) == 1 and req.num_sets > 1:
        diffs = [req.difficulty_levels[0]] * req.num_sets
    elif len(req.difficulty_levels) == req.num_sets:
        diffs = req.difficulty_levels
    else:
        raise ValueError("difficulty_levels must be length 1 (to broadcast) or equal to num_sets")

    wid = _id()
    sets_out = []

    for i, diff in enumerate(diffs, start=1):
        # ✅ generate dicts
        items_dicts = generate_items(context, req, diff, req.questions_per_set)

        # ✅ convert dicts → WorksheetItem models
        items = [WorksheetItem(**x) for x in items_dicts]

        worksheet_id = f"{_id()}-{diff}-set{i}"
        pdf_path = render_pdf(items, worksheet_id, req.target_language.lower())
        pdf_name = os.path.basename(pdf_path)

        sets_out.append(WorksheetSetResponse(
            set_no=i,
            difficulty=diff,  # type: ignore
            items=items,
            printable_pdf_url=f"/files/{pdf_name}"
        ))

    return WorksheetResponse(worksheet_id=wid, sets=sets_out)
